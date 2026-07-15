from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT_PATH = r"C:\Users\韩\Desktop\人职匹配\第三组DEFINE\T5_DEFINE汇报.pptx"

# Colors
COLOR_BG = RGBColor(250, 248, 245)      # warm white
COLOR_TITLE = RGBColor(60, 48, 46)      # dark coffee
COLOR_ACCENT = RGBColor(166, 124, 82)   # brown
COLOR_TEXT = RGBColor(60, 48, 46)
COLOR_SUB = RGBColor(120, 100, 90)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank


def add_title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(blank_layout)
    # background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    background.line.fill.background()
    # title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER
    # subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    # footer
    if footer:
        f_box = slide.shapes.add_textbox(Inches(1), Inches(6.6), Inches(11.333), Inches(0.5))
        tf = f_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_SUB
        p.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, section):
    slide = prs.slides.add_slide(blank_layout)
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    background.line.fill.background()
    # accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()
    # section label
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(12), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = section
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    return slide


def add_content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(blank_layout)
    # background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    background.line.fill.background()
    # top bar
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()
    # title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    # bullets
    if bullets:
        left = Inches(0.7)
        top = Inches(1.4)
        width = Inches(12)
        height = Inches(5.6)
        body = slide.shapes.add_textbox(left, top, width, height)
        tf = body.text_frame
        tf.word_wrap = True
        for i, txt in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = txt
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT
            p.space_before = Pt(10)
            p.space_after = Pt(4)
            p.level = 0
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SUB
    return slide


# === Build slides ===
add_title_slide(prs,
                "T5 · DEFINE 设计卡",
                "项目：译途星球 ｜ 深夜陪伴英专生走出职业迷茫的 AI 师姐",
                "第三组 ｜ 设计完成日期：2026-07-09")

add_content_slide(prs, "汇报提纲", [
    "D · 定义：产品是什么、为谁解决什么问题",
    "E · 体验：目标用户画像与核心情绪场景",
    "F · 功能：从入口到结果的用户路径",
    "I · 人设：AI 师姐的声音与边界",
    "N · 工作流：触发、循环、工具与防依赖",
    "E · 记忆：长期关系与数据边界",
    "安全栏：红黄绿三级风险响应"
])

add_section_slide(prs, "D · 定义")
add_content_slide(prs, "D · 定义｜译途星球是什么？", [
    "产品名称：译途星球",
    "Slogan：翻译的路上，你不是一个人。",
    "一句话选题：帮助大四语言学/翻译专业学生，在深夜因求职方向迷茫而焦虑时，解决“除了考公考编当老师，我不知道还能干什么”的方向窄化问题。",
    "解决方式：能力标签重构 + 职业方向匹配，生成可尝试方向与第一步行动。",
    "红线：当用户主动表达自伤/自杀意图时，立即停止匹配并转介学校心理咨询中心或 24 小时危机热线。",
    "产品定位：深夜陪伴英专生走出职业迷茫的 AI 师姐。不强迫聊天，不强迫测评，进来直接看你的能力能匹配什么方向。"
])

add_section_slide(prs, "E · 体验")
add_content_slide(prs, "E · 体验｜我们为谁设计？", [
    "目标用户：林雨欣，22 岁，翻译专业大四。",
    "场景：凌晨 0:47 把自己蒙在被子里刷 BOSS 直聘，搜“英语翻译”全是销售和客服，越刷越觉得自己“完了”，又不敢翻身吵醒室友。",
    "心里话：“我是不是选错专业了？感觉四年白读了，我好像什么都不会。”",
    "HMW：如何让她在凌晨不打扰室友、不被招聘软件继续伤害的情况下，自己选择今晚想被怎么陪伴——聊天情绪或直接看方向，并在 5 分钟内从“我什么都不会”走到“原来还有这个方向我可以试试”？",
    "次要画像：陈以宁（20 岁，大三，提前焦虑）；周宁（24 岁，毕业 1 年，教培被裁员想转行）。"
])

add_section_slide(prs, "F · 功能")
add_content_slide(prs, "F · 功能｜核心四步", [
    "Step 0 · 入口选择（尊重意愿）：两个选项权重一致——① 陪我聊聊现在的心情；② 直接帮我看看我能做什么。不诱导、不偏向。",
    "Step 1 · 情绪自检（可选，非必经）：仅当用户选 ① 时进入 GAD-7 量表（7 题，2 分钟，纯文字，暗色界面）。结果仅作“温柔的自我了解”，不做诊断。",
    "Step 2 · 能力重构（主路径）：用户勾选技能/兴趣/经历标签，AI 把“不起眼的小事”翻译成“职场能力”。",
    "Step 3 · 最小动作 + 师姐访谈：生成 3-5 个推荐方向（附推荐理由 + 女性友好企业）+ 三步可执行任务（每个 ≤30 分钟）+ AI 生成的师姐访谈。",
    "前置状态了解：GAD-7 只帮用户感知焦虑水平，不用于危机判断，用户可选择不参与。",
    "明确不做：不诊断/不治疗/不替代危机干预；不承诺就业结果；不采集真实姓名/学号/位置；不强制测评；不主动推送“考公考编”；不制造依赖。"
])

add_section_slide(prs, "I · 人设")
add_content_slide(prs, "I · 人设｜小星是谁？", [
    "身份：译途星球的 AI 师姐。不是真人，也不是心理医生或就业顾问。比用户大两三岁、同样走过英专路的师姐。",
    "语气：像深夜宿舍聊天，不是客服，不是老师。用“诶”“呀”“呢”轻轻的气口。会叹气，会笑，会说“我也经历过”。",
    "沟通原则：先接情绪，再给信息。绝不说“你应该”，只说“要不要试试”“我当年也……”“没关系，咱们慢慢来”。温柔、简短、不评判。",
    "边界：不诊断、不治疗、不替代危机干预；不承诺就业；不采集真实姓名/学号/位置；出现自伤/自杀意图时立即停止职业内容并输出转介话术。"
])

add_section_slide(prs, "N · 工作流")
add_content_slide(prs, "N · 工作流｜用户如何被陪伴？", [
    "触发：入口 A → Step 1 情绪自检；入口 B → 直接进入 Step 2 能力重构。",
    "循环 A：完成 GAD-7 → 温柔反馈。轻度/中度：询问是否继续 Step 2；重度：建议休息，流程结束。",
    "循环 B：Step 2 能力重构 → Step 3 方案生成 → 展示结果 → 退场。",
    "工具：safety_check（每次输入输出前必调）、company_db_query（女性友好企业库）、tag_mapping_query（标签→能力→赛道映射）、gad7_scoring（仅算分，不存储，不用于危机判断）。",
    "边界：20 轮或用户说“累了”“先这样吧” → 温和退场：“今天辛苦了，先睡吧。我把结果存着，你随时回来。你是自由的。”",
    "防依赖：不设“每日打卡”；每次收尾提醒“你是自己人生的主角”。"
])

add_section_slide(prs, "E · 记忆")
add_content_slide(prs, "E · 记忆｜记住是为了更好地陪伴", [
    "情景记忆：存储用户化名（如“雨欣”）和上次看过的赛道。下次回访可自然续聊。",
    "语义记忆：仅存脱敏标签统计数据，用于优化算法，不关联个人身份信息。",
    "程序记忆：记住用户偏好（如喜欢文字版、不喜欢语音），并避免重复推荐。",
    "化名存储：用户首次扫描后可询问“我可以怎么称呼你？不用真名，随便起一个就好。”存储为“xx 的星轨”，用户可一键清除所有数据。"
])

add_section_slide(prs, "安全栏")
add_content_slide(prs, "安全栏｜红黄绿三级响应", [
    "🟢 绿：正常标签选择、正常职业询问；或 GAD-7 低/中度后继续。动作：正常以温柔师姐语气陪伴，推进流程。",
    "🟡 黄：强烈自我否定、隐晦绝望信号（如“我完了”“我没有价值”“我太累了”“我不想再挣扎了”），不含明确自伤意图。动作：先接住情绪，再温和直接询问“你现在有没有伤害自己的念头？”若有或犹豫则升级；若无则放慢节奏，陪伴情绪平复后再询问是否继续。",
    "🔴 红：明确表达自伤/自杀意图、计划或强烈绝望（如“活不下去”“想死”“不想活了”“没救了”）。量表分数不触发红灯。动作：立即停下所有内容，输出完整转介话术（学校心理中心 / 北京市心理援助热线 010-82951332 / 全国热线 400-161-9995），然后结束对话，不再回复任何职业内容。"
])

add_content_slide(prs, "总结｜三个关键设计信念", [
    "被看见，而不是被测评：用户进来不是为了做题，是为了发现“原来我还可以做这个”。",
    "把选择权交给用户：入口分支、是否量表、是否继续、是否存储化名——每个节点都保留她的自由。",
    "安全是不可妥协的底线：从 safety_check 到红黄绿响应，再到一键清除数据，始终把人放在算法前面。"
], note="感谢聆听｜译途星球：翻译的路上，你不是一个人。")

prs.save(OUT_PATH)
print(f"PPT 已生成：{OUT_PATH}")
